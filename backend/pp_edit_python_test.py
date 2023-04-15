from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Cm

# フォントのルール
style_format = {
    'type' : {
        'font' : 'メイリオ',
        'font-size' : Pt(60),
        'font-color' : RGBColor(255, 255, 255),
    },
    'date' : {
        'font' : 'メイリオ',
        'font-size' : Pt(40),
        'font-color' : RGBColor(64, 64, 64),
    },
    'department-presentor' : {
        'font' : 'メイリオ',
        'font-size' : Pt(24),
        'font-color' : RGBColor(64, 64, 64),
    },
    'department-content' : {
        'font' : 'メイリオ',
        'font-size' : Pt(20),
        'font-color' : RGBColor(64, 64, 64),
    },
    'title' : {
        'font' : 'メイリオ',
        'font-size' : Pt(44),
        'font-color' : RGBColor(64, 64, 64),
    },
    'content' : {
        'font' : 'メイリオ',
        'font-size' : Pt(30),
        'font-color' : RGBColor(64, 64, 64),
    },
}
# その他のルール
# *~~~* : 太字

contents = [
    {
        'title' : 'title1',
        'content' : 'content1',
    },
    {
        'title' : 'title2',
        'content' : 'content2',
    },
]


# 対象の文字列に、フォント、フォントサイズ、フォントの色を一括しているするためのコード
def apply_format(pptx_obj, target):
    pptx_obj.font.name = style_format[target]['font']
    pptx_obj.font.size = style_format[target]['font-size']
    pptx_obj.font.color.rgb = style_format[target]['font-color']


# 新規スライドの追加
def add_slide(prs):
    slide_layout_6 = prs.slide_layouts[6]
    prs.slides.add_slide(slide_layout_6)


# 連絡事項、部門報告、目次などの文字列が記載される部分の文字列を入力するための処理
def set_type(slide, text):
    width = Cm(29.21)
    height = Cm(3.68)
    left = Cm(0.89)
    top = Cm(0.8)
    textbox = slide.shapes.add_textbox(
        left=left,
        top=top,
        width=width,
        height=height,
    )
    textbox.text = text
    paragraph = textbox.text_frame.paragraphs[0]
    apply_format(paragraph, 'type')
    paragraph.font.bold = True


# 連絡事項の内容部分のコード
def set_content(slide, text):
    width = Cm(29.21)
    height = Cm(12.09)
    left = Cm(2.33)
    top = Cm(5.07)
    textbox = slide.shapes.add_textbox(
        left=left,
        top=top,
        width=width,
        height=height,
    )
    textbox.text = text

    paragraphs = textbox.text_frame.paragraphs
    apply_format(paragraphs[0], 'title')
    paragraphs[0].font.bold = True

    for paragraph in paragraphs[1:]:
        apply_format(paragraph, 'content')

    textbox.text_frame.word_wrap = True


# 連絡事項のタイトル部分
def set_title_bg(slide):
    width = Cm(33.87)
    height = Cm(3.68)
    left = Cm(0)
    top = Cm(0)
    fill_color = RGBColor(68, 114, 196)
    shape = slide.shapes.add_shape(
        autoshape_type_id=1,  # 1は四角形を表す
        left=left,
        top=top,
        width=width,
        height=height
    )
    shape.z_order = 0
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color


# 最初のスライドに日付を追加するための処理
def set_date_text(slides):
    start_slide = slides[0]
    date_shape = start_slide.shapes[1]
    data_paragraph = date_shape.text_frame.paragraphs[0]
    data_paragraph.text = '4月8日'
    apply_format(data_paragraph, 'date')
    data_paragraph.font.bold = True


# 連絡事項の追加制御
def add_info_share(prs, info_contents):
    '''
    prs: Presentationオブジェクトを渡す
    info_contents: [
        'title' : '~~~',
        'contents : '~~~'.
    ]
    info_contentは連絡事項に関する情報だけを渡すようにする
    '''
    slides = prs.slides
    for content in info_contents:
        add_slide(prs)
        slide_cnt = len(slides)
        slide = slides[slide_cnt-1]
        set_title_bg(slide)
        set_type(slide, '連絡事項')
        info_content = f'{content.title}\n{content.content}'
        set_content(slide, info_content)
        

# ---------------これから-------------------------

# 目次部分の内容の追加
def set_outline_content(slide):
    set_title_bg(slide)
    set_type(slide, '目次')


# 部門報告の内容追加
def set_department_content(slide):
    set_title_bg(slide)
    set_type('部門報告')


def add_header_on_start_slide(start_slide):
    width = Cm(29.21)
    height = Cm(3.68)
    left = Cm(0.89)
    top = Cm(0.8)
    textbox = start_slide.shapes.add_textbox(
        left=left,
        top=top,
        width=width,
        height=height,
    )
    textbox.text = '報告会'
    paragraph = textbox.text_frame.paragraphs[0]
    apply_format(paragraph, 'type')
    paragraph.font.bold = True

def add_start_slide(prs):
    '''
    prs : 空のpresentationオブジェクトを渡す
    '''
    add_slide(prs)
    start_slide = prs.slides[0]
    add_header_on_start_slide(start_slide)
    
# -----------------これから---------------------



# ファイルの操作に必要な変数の定義
datefmt = '20230407'
output_filename = f'{datefmt}.pptx'
input_filename = 'sample.pptx'
input_path = f'./inputs/{input_filename}'

prs = Presentation(input_path)

add_start_slide(prs)

prs.save('./outputs/output.pptx')