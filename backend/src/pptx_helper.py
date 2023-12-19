# -*- coding: utf-8 -*-

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Cm, Pt

# フォントのルール
style_format = {
    'header_on_start_slide': {
        'font': 'メイリオ',
        'font-size': Pt(138),
        'font-color': RGBColor(64, 64, 64),
    },
    'type': {
        'font': 'メイリオ',
        'font-size': Pt(60),
        'font-color': RGBColor(255, 255, 255),
    },
    'date_outline': {
        'font': 'メイリオ',
        'font-size': Pt(40),
        'font-color': RGBColor(64, 64, 64),
    },
    'department-title': {
        'font': 'メイリオ',
        'font-size': Pt(32),
        'font-color': RGBColor(64, 64, 64),
    },
    'department-content': {
        'font': 'メイリオ',
        'font-size': Pt(20),
        'font-color': RGBColor(64, 64, 64),
    },
    'title': {
        'font': 'メイリオ',
        'font-size': Pt(44),
        'font-color': RGBColor(64, 64, 64),
    },
    'content': {
        'font': 'メイリオ',
        'font-size': Pt(30),
        'font-color': RGBColor(64, 64, 64),
    },
}

# その他のルール
# *~~~* : 太字


# 対象の文字列に、フォント、フォントサイズ、フォントの色を一括しているするためのコード
def apply_format(pptx_obj, target):
    pptx_obj.font.name = style_format[target]['font']
    pptx_obj.font.size = style_format[target]['font-size']
    pptx_obj.font.color.rgb = style_format[target]['font-color']


# 新規スライドの追加
def add_slide(prs):
    slide_layout_6 = prs.slide_layouts[6]
    prs.slides.add_slide(slide_layout_6)


# 連絡事項、部門報告、目次などの文字列が記載される部分の文字列を入力するための処理
def set_type(slide, text):
    width = Cm(29.21)
    height = Cm(3.68)
    left = Cm(0.89)
    top = Cm(0.8)
    textbox = slide.shapes.add_textbox(
        left=left,
        top=top,
        width=width,
        height=height,
    )
    textbox.text = text
    paragraph = textbox.text_frame.paragraphs[0]
    apply_format(paragraph, 'type')
    paragraph.font.bold = True


# 連絡事項の内容部分のコード
def set_content(slide, text):
    width = Cm(29.21)
    height = Cm(12.09)
    left = Cm(2.33)
    top = Cm(5.07)
    textbox = slide.shapes.add_textbox(
        left=left,
        top=top,
        width=width,
        height=height,
    )
    textbox.text = text

    paragraphs = textbox.text_frame.paragraphs
    apply_format(paragraphs[0], 'title')
    paragraphs[0].font.bold = True

    for paragraph in paragraphs[1:]:
        apply_format(paragraph, 'content')

    textbox.text_frame.word_wrap = True


# 連絡事項のタイトル部分
def set_title_bg(slide):
    width = Cm(33.87)
    height = Cm(3.68)
    left = Cm(0)
    top = Cm(0)
    fill_color = RGBColor(68, 114, 196)
    shape = slide.shapes.add_shape(
        autoshape_type_id=1,  # 1は四角形を表す
        left=left,
        top=top,
        width=width,
        height=height)
    shape.z_order = 0
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color


# 最初のスライドに日付を追加するための処理
def set_date_text(start_slide, datefmt: str):
    '''
    start_slide: 最初のスライド
    datefmt: 定例会の日時(例: 2023年4月7日（金）)
    '''
    width = Cm(25.4)
    height = Cm(1.98)
    left = Cm(4.23)
    top = Cm(10.42)
    textbox = start_slide.shapes.add_textbox(
        left=left,
        top=top,
        width=width,
        height=height,
    )
    text_frame = textbox.text_frame
    text_frame.text = datefmt
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    apply_format(paragraph, 'date_outline')
    paragraph.font.bold = True


# 連絡事項の追加制御
def add_info_share(prs, info_contents):
    '''
    prs: Presentationオブジェクトを渡す
    info_contents: [
        'title' : '~~',
        'contents : '~~'.
    ]
    info_contentは連絡事項に関する情報だけを渡すようにする
    '''
    slides = prs.slides
    for content in info_contents:
        add_slide(prs)
        slide_cnt = len(slides)
        slide = slides[slide_cnt - 1]
        set_title_bg(slide)
        set_type(slide, '連絡事項')
        info_content = '{}\n{}'.format(content['title'], content['content'])
        set_content(slide, info_content)


def add_header_on_start_slide(start_slide):
    width = Cm(25.4)
    height = Cm(6.63)
    left = Cm(4.23)
    top = Cm(4.53)
    textbox = start_slide.shapes.add_textbox(
        left=left,
        top=top,
        width=width,
        height=height,
    )
    text_frame = textbox.text_frame
    text_frame.text = '報告会'
    paragraph = textbox.text_frame.paragraphs[0]
    apply_format(paragraph, 'header_on_start_slide')
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.bold = True


def add_start_slide(prs, datefmt: str):
    '''
    prs : 空のpresentationオブジェクトを渡す
    datefmt : 定例会日時(例: 2023年4月7日（金）)
    '''
    add_slide(prs)
    start_slide = prs.slides[0]
    add_header_on_start_slide(start_slide)
    set_date_text(start_slide, datefmt)


# 目次部分の内容の追加
def add_outline_slide(prs, hasInfo: bool):
    add_slide(prs)
    slide = prs.slides[1]
    set_title_bg(slide)
    set_type(slide, '目次')
    set_outline_content(slide, hasInfo)


def set_outline_content(slide, hasInfo):
    width = Cm(31.73)
    height = Cm(13.67)
    left = Cm(0.97)
    top = Cm(4.52)
    textbox = slide.shapes.add_textbox(
        left=left,
        top=top,
        width=width,
        height=height,
    )
    text_frame = textbox.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.text = '1. 部門報告'
    paragraph.level = 0
    paragraph.font.bold = True
    apply_format(paragraph, 'date_outline')
    paragraph.line_spacing = 1.2

    if hasInfo:
        paragraph = text_frame.add_paragraph()
        paragraph.text = '2. 連絡事項'
        paragraph.level = 0
        paragraph.font.bold = True
        apply_format(paragraph, 'date_outline')
        paragraph.line_spacing = 1.2

    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.vertical_anchor = MSO_ANCHOR.TOP


# 部門報告の内容追加
def set_department_content(department_slide, departments_contents):
    '''
    department_slide: slideオブジェクト
        部門報告のスライド
    departments_contents: 各部門ごとの活動内容に関する辞書
        {
            'ds' : [
                {'title' : '~~', 'content', '~~'},
            ], ...
        }
    '''
    width = Cm(30.2)
    height = Cm(13.67)
    left = Cm(1.81)
    top = Cm(3.52)
    textbox = department_slide.shapes.add_textbox(
        left=left,
        top=top,
        width=width,
        height=height,
    )
    text_frame = textbox.text_frame
    departments = ['ds', 'de', 'biz', 'cc']
    for department in departments:
        department_upper = department.upper() if department != 'biz' else 'Biz'
        department_label = f'{department_upper}部門'
        department_title_paragraph = text_frame.add_paragraph()
        department_title_paragraph.text = f'{department_label}'
        apply_format(department_title_paragraph, 'department-title')
        department_title_paragraph.font.bold = True
        department_title_paragraph.line_spacing = 1.2

        if len(departments_contents[department]) == 0:
            department_content_paragraph = text_frame.add_paragraph()
            department_content_paragraph.line_spacing = 1.2
            department_content_paragraph.text = f'\t・なし'
            apply_format(department_content_paragraph, 'department-content')
            department_content_paragraph.font.bold = True
        else:
            for department_content in departments_contents[department]:
                department_content_paragraph = text_frame.add_paragraph()
                department_content_paragraph.line_spacing = 1.25
                department_content_paragraph.text = f'\t・{department_content}'
                apply_format(department_content_paragraph,
                             'department-content')
                department_content_paragraph.font.bold = True


def add_department_slide(prs, departments_contents):
    '''
    prs : プレゼンテーションオブジェクト
    departments_contents: 各部門ごとの活動内容に関する辞書
        {
            'ds' : [
                {'title' : '~~', 'content', '~~'},
            ], ...
        }
    '''
    add_slide(prs)
    slides = prs.slides
    department_slide = slides[2]
    set_title_bg(department_slide)
    set_type(department_slide, '１．各部門報告/その他報告')
    set_department_content(department_slide, departments_contents)
